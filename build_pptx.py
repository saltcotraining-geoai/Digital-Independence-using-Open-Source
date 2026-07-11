#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── colour palette ──
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE_ACCENT = RGBColor(0xFF, 0x8C, 0x00)

def set_slide_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_frame(slide, left, top, width, height, items, font_size=18, color=WHITE, bullet_char="\u25B8 "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}{item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf

def episode_summary_slide(num, title, concepts, deliverable):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # accent bar top
    add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    # number badge
    badge = add_shape_bg(slide, ACCENT, Inches(0.6), Inches(0.4), Inches(1.2), Inches(1.2))
    badge.text_frame.paragraphs[0].text = str(num)
    badge.text_frame.paragraphs[0].font.size = Pt(44)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = DARK
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.paragraphs[0].font.name = "Calibri"
    badge.text_frame.word_wrap = False
    # title
    add_textbox(slide, Inches(2.2), Inches(0.4), Inches(10), Inches(0.9), f"Episode {num}: {title}", font_size=36, bold=True, color=ACCENT)
    # details
    add_textbox(slide, Inches(2.2), Inches(1.3), Inches(10), Inches(0.5), f"Duration: 20 min  |  1.5+17+2 format", font_size=16, color=LIGHT_GREY)
    # concepts
    add_textbox(slide, Inches(0.6), Inches(2.2), Inches(5.5), Inches(0.5), "KEY CONCEPTS", font_size=14, bold=True, color=ORANGE_ACCENT)
    add_bullet_frame(slide, Inches(0.6), Inches(2.8), Inches(5.5), Inches(3.5), concepts, font_size=18)
    # deliverable box
    box = add_shape_bg(slide, RGBColor(0x22, 0x22, 0x3E), Inches(7), Inches(2.2), Inches(5.5), Inches(2))
    add_textbox(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(0.4), "HANDS-ON DELIVERABLE", font_size=14, bold=True, color=ORANGE_ACCENT)
    add_textbox(slide, Inches(7.3), Inches(3), Inches(5), Inches(1), deliverable, font_size=20, color=WHITE)
    # challenge
    add_shape_bg(slide, RGBColor(0x22, 0x22, 0x3E), Inches(7), Inches(4.6), Inches(5.5), Inches(2))
    add_textbox(slide, Inches(7.3), Inches(4.8), Inches(5), Inches(0.4), "HOMEWORK CHALLENGE", font_size=14, bold=True, color=ORANGE_ACCENT)
    return slide

def make_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
                "Digital Independence\nusing Open Source", font_size=48, bold=True, color=WHITE)
    add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                "Trainer Guide  \u2014  Level 1: Beginner Linux Course", font_size=26, color=ACCENT)
    add_textbox(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
                "A zero-to-expert curriculum for low-resource classrooms", font_size=18, color=LIGHT_GREY)
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
                "7 Episodes  \u25CF  20 minutes each  \u25CF  Offline-first  \u25CF  Free & open source",
                font_size=16, color=LIGHT_GREY)
    return slide

def make_section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_bg(slide, ACCENT, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1), title, font_size=40, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6), subtitle, font_size=20, color=LIGHT_GREY)
    return slide

# ══════════════════════════════════════
#   SLIDE 1 – TITLE
# ══════════════════════════════════════
make_title_slide()

# ══════════════════════════════════════
#   SLIDE 2 – MISSION
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Our Mission", font_size=36, bold=True, color=ACCENT)
mission_items = [
    "Open source is our vehicle to share knowledge, build skills, and teach those with limited resources to participate and prosper.",
    "This curriculum is free, public, and forever free. Licensed under CC BY-SA.",
    "Every student with an old laptop, unstable internet, or no money for licenses is exactly who this serves.",
    "Teachers are encouraged to copy, modify, translate, and teach this anywhere in the world.",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(11.5), Inches(4), mission_items, font_size=22)
add_shape_bg(slide, RGBColor(0x22, 0x22, 0x3E), Inches(0.6), Inches(5), Inches(12), Inches(1.5))
add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.2),
            "Special thanks to: Linux Mint & XFCE teams  \u25CF  Debian project & Ubuntu community  \u25CF  GNU project  \u25CF  Every open-source maintainer",
            font_size=16, color=ORANGE_ACCENT)

# ══════════════════════════════════════
#   SLIDE 3 – CAREER PIPELINE
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Career Pipeline", font_size=36, bold=True, color=ACCENT)
levels = [
    ("PREREQUISITE DOORWAY", "Digital Independence\n(Linux CLI Fundamentals)", "You Are Here"),
    ("LEVEL 1", "Desktop Mapping\nQGIS, Vector & Raster", "Coming Next"),
    ("LEVEL 2", "Python Automation\nGeoPandas, Scripting", "Coming Next"),
    ("LEVEL 3", "Web Geography\nFolium, OpenLayers", "Coming Next"),
]
y_start = 1.6
for i, (level, desc, status) in enumerate(levels):
    y = y_start + i * 1.4
    color = ACCENT if i == 0 else LIGHT_GREY
    box = add_shape_bg(slide, RGBColor(0x22, 0x22, 0x3E), Inches(0.6), Inches(y), Inches(12), Inches(1.1))
    add_textbox(slide, Inches(0.9), Inches(y + 0.1), Inches(3), Inches(0.9), level, font_size=16, bold=True, color=color)
    add_textbox(slide, Inches(4.2), Inches(y + 0.1), Inches(5), Inches(0.9), desc, font_size=18, color=WHITE)
    add_textbox(slide, Inches(10), Inches(y + 0.1), Inches(2.5), Inches(0.9), status, font_size=14, color=ORANGE_ACCENT)
    if i < 3:
        add_textbox(slide, Inches(6.5), Inches(y + 0.9), Inches(0.5), Inches(0.5), "\u25BC", font_size=14, color=ACCENT)

# ══════════════════════════════════════
#   SLIDE 4 – RULES
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Course Rules", font_size=36, bold=True, color=ACCENT)
rules = [
    "Level 1 must be completed BEFORE any other level. No exceptions.",
    "Every episode is exactly 20 minutes: 1.5 min intro + 17 min hands-on + 1.5 min outro.",
    "Students must complete the homework challenge before moving to the next episode.",
    "All software, data, and materials must work 100% offline.",
    "Terminal font must be 16pt+ with high contrast for low-resolution video viewers.",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5), rules, font_size=22)

# ══════════════════════════════════════
#   SLIDE 5 – HARDWARE
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Hardware Requirements", font_size=36, bold=True, color=ACCENT)
specs = [
    ("Component", "Minimum", "Recommended"),
    ("RAM", "2 GB", "4 GB"),
    ("Storage", "20 GB free", "100 GB free"),
    ("USB Drive", "8 GB", "16 GB (for persistence)"),
    ("OS", "Linux Mint XFCE", "Lubuntu / MX Linux"),
]
y = 1.6
for i, (col1, col2, col3) in enumerate(specs):
    bg = ACCENT if i == 0 else RGBColor(0x22, 0x22, 0x3E)
    box = add_shape_bg(slide, bg, Inches(0.6), Inches(y), Inches(12), Inches(0.6))
    add_textbox(slide, Inches(0.9), Inches(y + 0.05), Inches(4), Inches(0.5), col1,
                font_size=18, bold=(i == 0), color=DARK if i == 0 else WHITE)
    add_textbox(slide, Inches(5), Inches(y + 0.05), Inches(3), Inches(0.5), col2,
                font_size=18, bold=(i == 0), color=DARK if i == 0 else WHITE)
    add_textbox(slide, Inches(8.5), Inches(y + 0.05), Inches(3), Inches(0.5), col3,
                font_size=18, bold=(i == 0), color=DARK if i == 0 else WHITE)
    y += 0.7

add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
            "Tip: Enable USB persistence so student work survives reboots.\nUse Rufus (Windows) or Balena Etcher (Mac/Linux) to create bootable drives.",
            font_size=16, color=LIGHT_GREY)

# ══════════════════════════════════════
#   SLIDES 6-12 – EPISODES
# ══════════════════════════════════════
make_section_slide("Episode Overview", "7 episodes  \u25CF  20 min each  \u25CF  Level 1: Digital Independence")

episodes = [
    (1, "The Digital Rebirth",
     ["Why leave Windows/Mac?", "FOSS philosophy (the recipe analogy)", "free -h: check memory usage", "df -h: check storage usage"],
     "Boot Linux from USB. Run free -h and df -h. Record your baseline metrics."),
    (2, "The Digital Scout",
     ["pwd: where am I?", "ls / ls -F: what is here?", "cd: navigate folders", "Tab key autocomplete", "Maze game navigation"],
     "Navigate the Episode2_Maze using only cd and ls. Find the treasure.txt password."),
    (3, "The Digital Architect",
     ["mkdir: create folders", "touch: create empty files", "nano: terminal text editor", "cat / less: read files"],
     "Build a folder structure. Write your bio in nano. Save and exit using Ctrl+O / Ctrl+X."),
    (4, "The Digital Commander",
     ["cp: copy files safely", "mv: move / rename", "rm: delete (NO recycle bin)", "rm -r: delete folders"],
     "Backup report.txt to backups/. Rename to final_report.txt. Delete junk.tmp."),
    (5, "The Software Alchemist",
     ["sudo: administrator privileges", "dpkg -i: install local .deb files", "Offline package distribution"],
     "Install neofetch from a local .deb file using sudo dpkg -i. Run neofetch to verify."),
    (6, "The System Savior",
     ["htop / top: real-time process viewer", "PID: Process ID", "pkill: kill by name", "kill -9: nuclear force quit"],
     "Run htop, identify top 3 memory processes. Kill a harmless app with pkill."),
    (7, "Bonus: Automation Bridge",
     ["#!/bin/bash: the shebang", "chmod +x: make executable", "echo: print text", ">: redirect output to file"],
     "Write system_guard.sh. Run it. Modify it to save a report file."),
]
for num, title, concepts, deliverable in episodes:
    episode_summary_slide(num, title, concepts, deliverable)

# ══════════════════════════════════════
#   SLIDE – YOUTUBE RECORDING TIPS
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Recording Tips for YouTube", font_size=36, bold=True, color=ACCENT)
tips = [
    "Terminal font: 16pt or higher, monospace (Ubuntu Mono or Monospace Bold)",
    "Theme: White text on black background (minimizes glare on cheap screens)",
    "Zoom in to 150-200% so text is legible even at 360p playback",
    "Record at 1080p, export at 720p for small file sizes (under 100 MB per episode)",
    "Run clear before every demo section for a clean slate",
    "Add YouTube chapter markers in each video description",
    "Hide desktop icons and close all other windows before recording",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5), tips, font_size=20)

# ══════════════════════════════════════
#   SLIDE – OFFLINE SETUP
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Offline Classroom Setup", font_size=36, bold=True, color=ACCENT)
steps = [
    "Download Linux Mint XFCE ISO and create bootable USBs (enable persistence)",
    "Pre-download packages: sudo apt-get install --download-only neofetch htop",
    "Copy .deb files from /var/cache/apt/archives/ to each student's machine",
    "Run setup_script.sh to create the classroom folder with maze + sandbox",
    "Print laminated cheat sheets for pwd, ls, cd, cp, mv, rm, nano shortcuts",
    "Test everything offline before class begins",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5), steps, font_size=20)
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.8),
            "Pro tip: One internet-connected machine can serve an entire off-grid classroom via USB distribution.",
            font_size=16, color=ORANGE_ACCENT)

# ══════════════════════════════════════
#   SLIDE – BEFORE EACH EPISODE
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Episode Prerequisites Checklist", font_size=36, bold=True, color=ACCENT)
checks = [
    "Ep 1: Bootable Linux USB ready for each student",
    "Ep 2: Run setup_script.sh to create the Maze directories",
    "Ep 3: Episode 2 completed by all students",
    "Ep 4: Episode 3 completed by all students",
    "Ep 5: .deb package files copied to scripts/packages/",
    "Ep 6: htop installed (include in pre-downloaded packages)",
    "Bonus: check_homework.sh ready to grade final exam",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5), checks, font_size=20)

# ══════════════════════════════════════
#   FINAL SLIDE – CREDITS
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), "Thank You & Keep Building", font_size=36, bold=True, color=ACCENT)
credits = [
    "This curriculum is dedicated to the global open-source community.",
    "Thousands of developers gave their time for free so knowledge would not be locked behind paywalls.",
    "You are now part of that community. When you help someone else learn, you continue the chain.",
    "Licensed under CC BY-SA 4.0. Copy, modify, translate, teach — anywhere, any time, for free.",
    "",
    "saltcotraining-geoai  \u25CF  Digital Independence using Open Source",
]
add_bullet_frame(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(4), credits, font_size=22)

# ── SAVE ──
out_path = os.path.expanduser("~/Desktop/Digital_Indepedence_using_Open_Source/Trainer_Presentation.pptx")
prs.save(out_path)
print(f"Saved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
